#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX 幻灯片导出器

将生成的幻灯片序列导出为 PowerPoint (.pptx) 格式。
依赖: pip install python-pptx
"""

from pathlib import Path
from typing import Dict, Any, List, Optional


def _hex_to_rgb(hex_color: str):
    """将 #RRGGBB 转为 (r, g, b) 元组，供 python-pptx 使用"""
    hex_color = hex_color.strip()
    if hex_color.startswith('#'):
        hex_color = hex_color[1:]
    if len(hex_color) == 6:
        return (
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    return (0, 0, 0)


class PPTXExporter:
    """PPTX 幻灯片导出器（需安装 python-pptx）"""

    def __init__(self, width_inches: float = 13.333, height_inches: float = 7.5):
        """
        初始化导出器（16:9 默认尺寸）

        Args:
            width_inches: 幻灯片宽度（英寸）
            height_inches: 幻灯片高度（英寸）
        """
        self.width_inches = width_inches
        self.height_inches = height_inches

    def export(self, slides: List[Dict[str, Any]], metadata: Dict[str, Any]) -> bytes:
        """
        导出幻灯片为 PPTX 二进制内容

        Args:
            slides: 幻灯片数据列表（与 RevealJS 相同的 dict 结构）
            metadata: 元数据（title, author 等）

        Returns:
            bytes: PPTX 文件二进制内容
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
        except ImportError:
            raise ImportError(
                "导出 PPTX 需要安装 python-pptx: pip install python-pptx"
            )

        prs = Presentation()
        prs.slide_width = Inches(self.width_inches)
        prs.slide_height = Inches(self.height_inches)

        # 设置演示文稿属性
        core_props = prs.core_properties
        core_props.title = metadata.get('title', 'Presentation')
        core_props.author = metadata.get('author', '')

        for i, slide_data in enumerate(slides):
            # 第一张用标题版式，其余用标题和内容
            if i == 0:
                layout = prs.slide_layouts[6]  # 空白
            else:
                layout = prs.slide_layouts[6]  # 空白，便于自由放置

            slide = prs.slides.add_slide(layout)

            # 背景色
            bg_color = slide_data.get('background', {}).get('fillColor', '#FFFFFF')
            if bg_color:
                try:
                    r, g, b = _hex_to_rgb(bg_color)
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass

            # 标题（幻灯片标题栏）
            title_text = slide_data.get('title', '')
            if title_text:
                left = Inches(0.5)
                top = Inches(0.3)
                width = Inches(self.width_inches - 1)
                height = Inches(0.8)
                tx = slide.shapes.add_textbox(left, top, width, height)
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(28)
                p.font.bold = True
                try:
                    first_content = next(
                        (c for c in slide_data.get('contents', []) if c.get('type') == 'title'),
                        None
                    )
                    if first_content:
                        style = first_content.get('style', {})
                        color = style.get('color', '#1A1A1A')
                        r, g, b = _hex_to_rgb(color)
                        p.font.color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass

            # 内容区域
            contents = slide_data.get('contents', [])
            body_top = Inches(1.2)
            left = Inches(0.5)
            width = Inches(self.width_inches - 1)
            line_height = Inches(0.4)

            for content in contents:
                ctype = content.get('type', 'body')
                text = content.get('text', '')
                if not text:
                    continue

                if ctype == 'title':
                    # 已在上面画过主标题，这里可跳过或作为副标题
                    continue
                if ctype == 'body':
                    box = slide.shapes.add_textbox(left, body_top, width, Inches(1.2))
                    box.text_frame.word_wrap = True
                    box.text_frame.paragraphs[0].text = text
                    box.text_frame.paragraphs[0].font.size = Pt(18)
                    body_top += line_height + Inches(0.2)
                elif ctype == 'bullet':
                    level = content.get('level', 0)
                    indent = 0.3 + level * 0.3
                    box = slide.shapes.add_textbox(
                        Inches(0.5 + indent), body_top, Inches(self.width_inches - 1 - indent), Inches(0.5)
                    )
                    box.text_frame.word_wrap = True
                    p = box.text_frame.paragraphs[0]
                    p.text = "• " + text
                    p.font.size = Pt(20)
                    p.space_after = Pt(6)
                    body_top += line_height
                elif ctype == 'quote':
                    box = slide.shapes.add_textbox(left, body_top, width, Inches(1.0))
                    box.text_frame.word_wrap = True
                    p = box.text_frame.paragraphs[0]
                    p.text = '"' + text + '"'
                    p.font.italic = True
                    p.font.size = Pt(20)
                    body_top += Inches(0.8)
                elif ctype == 'cta_button':
                    box = slide.shapes.add_textbox(left, body_top, width, Inches(0.5))
                    p = box.text_frame.paragraphs[0]
                    p.text = text
                    p.font.bold = True
                    p.font.size = Pt(18)
                    body_top += Inches(0.6)

            # 备注
            notes_text = slide_data.get('notes', '')
            if notes_text and hasattr(slide, 'notes_slide'):
                try:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes_text
                except Exception:
                    pass

        # 写入内存
        import io
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def save(self,
             slides: List[Dict[str, Any]],
             metadata: Dict[str, Any],
             filepath: Path) -> Path:
        """
        保存为 .pptx 文件

        Args:
            slides: 幻灯片数据
            metadata: 元数据
            filepath: 输出路径

        Returns:
            Path: 保存后的文件路径
        """
        data = self.export(slides, metadata)
        path = Path(filepath)
        if path.suffix.lower() != '.pptx':
            path = path.with_suffix('.pptx')
        path.write_bytes(data)
        return path


def export_pptx(slides: List[Dict[str, Any]],
                metadata: Dict[str, Any],
                output_path: str,
                width_inches: float = 13.333,
                height_inches: float = 7.5) -> str:
    """
    便捷的 PPTX 导出函数

    Args:
        slides: 幻灯片数据
        metadata: 元数据
        output_path: 输出文件路径
        width_inches: 宽度（英寸）
        height_inches: 高度（英寸）

    Returns:
        str: 输出文件路径
    """
    exporter = PPTXExporter(width_inches=width_inches, height_inches=height_inches)
    return str(exporter.save(slides, metadata, Path(output_path)))


if __name__ == "__main__":
    import argparse
    import json

    parser = argparse.ArgumentParser(description="Export to PPTX")
    parser.add_argument("input", help="Input JSON file or 'demo'")
    parser.add_argument("--output", "-o", required=True, help="Output .pptx file")
    parser.add_argument("--width", type=float, default=13.333, help="Slide width (inches)")
    parser.add_argument("--height", type=float, default=7.5, help="Slide height (inches)")
    args = parser.parse_args()

    if args.input == 'demo':
        slides = [
            {
                'title': '欢迎',
                'background': {'fillColor': '#003366'},
                'transition': {'effect': 'fade'},
                'contents': [
                    {'type': 'title', 'text': '演示文稿标题', 'style': {'color': '#FFFFFF'}},
                    {'type': 'body', 'text': '副标题或描述', 'style': {'color': '#CCCCCC'}}
                ],
                'notes': ''
            },
            {
                'title': '内容',
                'background': {'fillColor': '#FFFFFF'},
                'transition': {'effect': 'slide'},
                'contents': [
                    {'type': 'title', 'text': '主要章节', 'style': {'color': '#003366'}},
                    {'type': 'bullet', 'text': '第一点', 'level': 0},
                    {'type': 'bullet', 'text': '第二点', 'level': 0},
                    {'type': 'bullet', 'text': '第三点', 'level': 0}
                ],
                'notes': ''
            }
        ]
        metadata = {'title': '演示文稿', 'author': '作者'}
    else:
        with open(args.input, 'r', encoding='utf-8') as f:
            data = json.load(f)
            slides = data.get('slides', [])
            meta = data.get('metadata', {})
            metadata = {
                'title': meta.get('title', 'Presentation'),
                'author': meta.get('author', '')
            }

    exporter = PPTXExporter(width_inches=args.width, height_inches=args.height)
    out = exporter.save(slides, metadata, Path(args.output))
    print(f"Saved to: {out}")
