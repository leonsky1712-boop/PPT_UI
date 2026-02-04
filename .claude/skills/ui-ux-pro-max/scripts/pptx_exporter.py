#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX 幻灯片导出器

将生成的幻灯片序列导出为 PowerPoint (.pptx) 格式。
支持：标题、副标题、正文、列表、引用、CTA、图片占位/本地图片、图表（柱状/折线/饼图）。
依赖: pip install python-pptx
"""

import io
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple


def _hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """将 #RRGGBB 转为 (r, g, b) 元组，供 python-pptx 使用"""
    hex_color = (hex_color or "").strip()
    if hex_color.startswith('#'):
        hex_color = hex_color[1:]
    if len(hex_color) == 6:
        return (
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    return (0, 0, 0)


class PPTXExporter:
    """PPTX 幻灯片导出器（需安装 python-pptx），支持图片与图表"""

    def __init__(self, width_inches: float = 13.333, height_inches: float = 7.5):
        """
        初始化导出器（16:9 默认尺寸）

        Args:
            width_inches: 幻灯片宽度（英寸）
            height_inches: 幻灯片高度（英寸）
        """
        self.width_inches = width_inches
        self.height_inches = height_inches

    def _ensure_pptx_imports(self):
        """延迟导入 python-pptx 并返回常用符号"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            return Presentation, Inches, Pt, RgbColor
        except ImportError:
            raise ImportError("导出 PPTX 需要安装 python-pptx: pip install python-pptx")

    def export(self, slides: List[Dict[str, Any]], metadata: Dict[str, Any]) -> bytes:
        """
        导出幻灯片为 PPTX 二进制内容

        Args:
            slides: 幻灯片数据列表（与 generator 的 to_dict 结构一致）
            metadata: 元数据（title, author, subtitle 等）

        Returns:
            bytes: PPTX 文件二进制内容
        """
        Presentation, Inches, Pt, RgbColor = self._ensure_pptx_imports()
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
        except ImportError:
            CategoryChartData = None
            XL_CHART_TYPE = None

        prs = Presentation()
        prs.slide_width = Inches(self.width_inches)
        prs.slide_height = Inches(self.height_inches)
        core_props = prs.core_properties
        core_props.title = metadata.get('title', 'Presentation')
        core_props.author = metadata.get('author', '')

        for i, slide_data in enumerate(slides):
            layout = prs.slide_layouts[6]  # 空白
            slide = prs.slides.add_slide(layout)
            bg = slide_data.get('background') or {}
            bg_color = bg.get('fillColor') or bg.get('fill_color', '#FFFFFF')
            if bg_color:
                try:
                    r, g, b = _hex_to_rgb(bg_color)
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass

            slide_type = slide_data.get('type', 'content')
            contents = slide_data.get('contents', [])
            title_text = slide_data.get('title', '')

            # 标题页 / 节分隔页：大标题 + 副标题
            if i == 0 or slide_type == 'title' or slide_type == 'section_divider':
                self._add_title_slide_content(slide, slide_data, contents, title_text, Inches, Pt, RgbColor)
            else:
                self._add_standard_slide_content(
                    slide, slide_data, contents, title_text,
                    Inches, Pt, RgbColor, CategoryChartData, XL_CHART_TYPE
                )

            notes_text = slide_data.get('notes', '')
            if notes_text:
                try:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes_text
                except Exception:
                    pass

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def _add_title_slide_content(self, slide, slide_data, contents, title_text, Inches, Pt, RgbColor):
        """标题页：居中大标题 + 副标题"""
        left = Inches(0.5)
        width = Inches(self.width_inches - 1)
        top = Inches(1.8)
        # 主标题（可用 contents 里 type=title 的 text，或 slide_data.title）
        main_title = title_text
        title_content = next((c for c in contents if c.get('type') == 'title'), None)
        if title_content and title_content.get('text'):
            main_title = title_content.get('text', main_title)
        if main_title:
            tx = slide.shapes.add_textbox(left, top, width, Inches(1.2))
            tf = tx.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = main_title
            p.font.size = Pt(40)
            p.font.bold = True
            style = title_content.get('style', {}) if title_content else {}
            color = style.get('color', '#1A1A1A')
            try:
                r, g, b = _hex_to_rgb(color)
                p.font.color.rgb = RgbColor(r, g, b)
            except Exception:
                pass
            top += Inches(1.4)
        subtitle_content = next((c for c in contents if c.get('type') == 'subtitle'), None)
        if subtitle_content and subtitle_content.get('text'):
            sub_tx = slide.shapes.add_textbox(left, top, width, Inches(0.8))
            sub_tx.text_frame.word_wrap = True
            sp = sub_tx.text_frame.paragraphs[0]
            sp.text = subtitle_content['text']
            sp.font.size = Pt(24)
            style = subtitle_content.get('style', {})
            color = style.get('color', '#444444')
            try:
                r, g, b = _hex_to_rgb(color)
                sp.font.color.rgb = RgbColor(r, g, b)
            except Exception:
                pass

    def _add_standard_slide_content(self, slide, slide_data, contents, title_text,
                                   Inches, Pt, RgbColor, CategoryChartData, XL_CHART_TYPE):
        """普通页：顶部标题 + 正文/列表/引用/图片/图表"""
        left = Inches(0.5)
        width = Inches(self.width_inches - 1)
        body_top = Inches(0.35)
        if title_text:
            tx = slide.shapes.add_textbox(left, body_top, width, Inches(0.65))
            tx.text_frame.word_wrap = True
            p = tx.text_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(28)
            p.font.bold = True
            first_title = next((c for c in contents if c.get('type') == 'title'), None)
            if first_title:
                style = first_title.get('style', {})
                try:
                    r, g, b = _hex_to_rgb(style.get('color', '#1A1A1A'))
                    p.font.color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass
            body_top = Inches(1.15)
        line_height = Inches(0.38)

        for content in contents:
            ctype = content.get('type', 'body')
            text = content.get('text', '')
            style = content.get('style', {})

            if ctype == 'title':
                continue
            if ctype == 'subtitle':
                box = slide.shapes.add_textbox(left, body_top, width, Inches(0.5))
                box.text_frame.word_wrap = True
                bp = box.text_frame.paragraphs[0]
                bp.text = text
                bp.font.size = Pt(20)
                try:
                    r, g, b = _hex_to_rgb(style.get('color', '#444444'))
                    bp.font.color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass
                body_top += Inches(0.55)
                continue
            if ctype == 'body':
                box = slide.shapes.add_textbox(left, body_top, width, Inches(1.0))
                box.text_frame.word_wrap = True
                bp = box.text_frame.paragraphs[0]
                bp.text = text
                bp.font.size = Pt(18)
                try:
                    r, g, b = _hex_to_rgb(style.get('color', '#333333'))
                    bp.font.color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass
                body_top += line_height + Inches(0.15)
                continue
            if ctype == 'bullet':
                level = content.get('level', 0)
                indent = 0.3 + level * 0.25
                box = slide.shapes.add_textbox(
                    Inches(0.5 + indent), body_top, Inches(self.width_inches - 1 - indent), Inches(0.5)
                )
                box.text_frame.word_wrap = True
                bp = box.text_frame.paragraphs[0]
                bp.text = "• " + text
                bp.font.size = Pt(20)
                bp.space_after = Pt(6)
                try:
                    r, g, b = _hex_to_rgb(style.get('color', '#333333'))
                    bp.font.color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass
                body_top += line_height
                continue
            if ctype == 'quote':
                box = slide.shapes.add_textbox(left, body_top, width, Inches(1.0))
                box.text_frame.word_wrap = True
                bp = box.text_frame.paragraphs[0]
                bp.text = '"' + text + '"'
                bp.font.italic = True
                bp.font.size = Pt(20)
                try:
                    r, g, b = _hex_to_rgb(style.get('color', '#555555'))
                    bp.font.color.rgb = RgbColor(r, g, b)
                except Exception:
                    pass
                body_top += Inches(0.9)
                continue
            if ctype == 'cta_button':
                box = slide.shapes.add_textbox(left, body_top, width, Inches(0.5))
                bp = box.text_frame.paragraphs[0]
                bp.text = text
                bp.font.bold = True
                bp.font.size = Pt(18)
                body_top += Inches(0.55)
                continue
            if ctype == 'image':
                body_top = self._add_image_content(slide, content, left, body_top, width, Inches, Pt, RgbColor)
                continue
            if ctype == 'chart' and CategoryChartData is not None and XL_CHART_TYPE is not None:
                body_top = self._add_chart_content(slide, content, left, body_top, Inches, CategoryChartData, XL_CHART_TYPE)
                continue
            if text:
                box = slide.shapes.add_textbox(left, body_top, width, Inches(0.5))
                box.text_frame.paragraphs[0].text = text
                box.text_frame.paragraphs[0].font.size = Pt(18)
                body_top += line_height

    def _add_image_content(self, slide, content: Dict, left, top, width, Inches, Pt, RgbColor) -> float:
        """添加图片：本地路径则插入图片，否则插入占位说明。left/width 为 Inches，top 为 float。返回下一行 top。"""
        image_path = content.get('path') or content.get('src') or content.get('url') or ''
        if isinstance(image_path, str) and image_path.startswith('file://'):
            image_path = image_path[7:]
        path = Path(image_path) if image_path else None
        try:
            if path and path.exists() and path.is_file():
                slide.shapes.add_picture(str(path), left, Inches(top), width=width, height=Inches(3.0))
                return top + 3.25
        except Exception:
            pass
        placeholder_text = content.get('text') or ('图片: ' + image_path if image_path else '图片占位符')
        box = slide.shapes.add_textbox(left, Inches(top), width, Inches(1.2))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = placeholder_text
        p.font.size = Pt(14)
        p.font.italic = True
        try:
            r, g, b = _hex_to_rgb('#888888')
            p.font.color.rgb = RgbColor(r, g, b)
        except Exception:
            pass
        return top + 1.35

    def _add_chart_content(self, slide, content: Dict, left, top, Inches,
                          CategoryChartData, XL_CHART_TYPE) -> float:
        """添加图表。content 可含 chart_type: bar|line|pie, data/categories。left 为 Inches，top 为 float。返回下一行 top。"""
        chart_type = (content.get('chart_type') or content.get('chartType') or 'bar').lower()
        data = content.get('data') or content.get('values') or [10, 20, 30, 40]
        categories = content.get('categories') or ['A', 'B', 'C', 'D'][:len(data)]
        series_name = content.get('series_name') or content.get('seriesName') or '数据'
        if isinstance(data, dict):
            categories = list(data.keys())
            data = list(data.values())
        if not categories or len(categories) != len(data):
            categories = [str(i) for i in range(len(data))]
        try:
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(series_name, tuple(float(x) for x in data))
            x, y = left, Inches(top)
            cx, cy = Inches(10), Inches(3.5)
            type_map = {
                'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE,
            }
            pptx_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            slide.shapes.add_chart(pptx_type, x, y, cx, cy, chart_data)
            return top + 3.8
        except Exception:
            return top + 0.5

    def save(self,
             slides: List[Dict[str, Any]],
             metadata: Dict[str, Any],
             filepath: Path) -> Path:
        """
        保存为 .pptx 文件

        Args:
            slides: 幻灯片数据
            metadata: 元数据
            filepath: 输出路径

        Returns:
            Path: 保存后的文件路径
        """
        data = self.export(slides, metadata)
        path = Path(filepath)
        if path.suffix.lower() != '.pptx':
            path = path.with_suffix('.pptx')
        path.write_bytes(data)
        return path


def export_pptx(slides: List[Dict[str, Any]],
                metadata: Dict[str, Any],
                output_path: str,
                width_inches: float = 13.333,
                height_inches: float = 7.5) -> str:
    """
    便捷的 PPTX 导出函数

    Args:
        slides: 幻灯片数据
        metadata: 元数据
        output_path: 输出文件路径
        width_inches: 宽度（英寸）
        height_inches: 高度（英寸）

    Returns:
        str: 输出文件路径
    """
    exporter = PPTXExporter(width_inches=width_inches, height_inches=height_inches)
    return str(exporter.save(slides, metadata, Path(output_path)))


if __name__ == "__main__":
    import argparse
    import json

    parser = argparse.ArgumentParser(description="Export to PPTX")
    parser.add_argument("input", help="Input JSON file or 'demo'")
    parser.add_argument("--output", "-o", required=True, help="Output .pptx file")
    parser.add_argument("--width", type=float, default=13.333, help="Slide width (inches)")
    parser.add_argument("--height", type=float, default=7.5, help="Slide height (inches)")
    args = parser.parse_args()

    if args.input == 'demo':
        slides = [
            {
                'title': '欢迎',
                'background': {'fillColor': '#003366'},
                'transition': {'effect': 'fade'},
                'contents': [
                    {'type': 'title', 'text': '演示文稿标题', 'style': {'color': '#FFFFFF'}},
                    {'type': 'body', 'text': '副标题或描述', 'style': {'color': '#CCCCCC'}}
                ],
                'notes': ''
            },
            {
                'title': '内容',
                'background': {'fillColor': '#FFFFFF'},
                'transition': {'effect': 'slide'},
                'contents': [
                    {'type': 'title', 'text': '主要章节', 'style': {'color': '#003366'}},
                    {'type': 'bullet', 'text': '第一点', 'level': 0},
                    {'type': 'bullet', 'text': '第二点', 'level': 0},
                    {'type': 'bullet', 'text': '第三点', 'level': 0}
                ],
                'notes': ''
            }
        ]
        metadata = {'title': '演示文稿', 'author': '作者'}
    else:
        with open(args.input, 'r', encoding='utf-8') as f:
            data = json.load(f)
            slides = data.get('slides', [])
            meta = data.get('metadata', {})
            metadata = {
                'title': meta.get('title', 'Presentation'),
                'author': meta.get('author', '')
            }

    exporter = PPTXExporter(width_inches=args.width, height_inches=args.height)
    out = exporter.save(slides, metadata, Path(args.output))
    print(f"Saved to: {out}")
